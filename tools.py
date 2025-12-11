# tools.py
from typing import Dict, Optional
import pandas as pd
import speech_recognition as sr
# from pydx2 import PdfReader
import docx2txt
import markdown
import re
from moviepy import VideoFileClip
from openpyxl import load_workbook
import json
from pptx import Presentation
import os
# from crewai_tools import PDFSearchTool
from PyPDF2 import PdfReader
from pathlib import Path
from datetime import datetime
import requests
from openai import OpenAI
from dotenv import load_dotenv
from ftplib import FTP
from urllib.parse import urlparse
import logging
import paramiko

load_dotenv()  # Load environment variables from .env file

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Airtable API credentials
AIRTABLE_API_TOKEN = os.getenv("AIRTABLE_API_TOKEN")
BASE_ID = os.getenv("BASE_ID")
TABLE_NAME = os.getenv("TABLE_NAME")
AIRTABLE_URL = f"https://api.airtable.com/v0/{BASE_ID}/{TABLE_NAME}"
# OpenAI API Key
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY")) # Replace with environment variable in production

PLATFORM_LIMITS = {
    "twitter": {"chars": 280, "words": None},
    "instagram": {"chars": None, "words": 400},
    "linkedin": {"chars": None, "words": 600},
    "facebook": {"chars": None, "words": 1000},
    "wordpress": {"chars": None, "words": 2000},
    "youtube": {"chars": None, "words": 2000},
    "tiktok": {"chars": None, "words": 400},
}

UPLOAD_DIR = './uploads'

# # Helper function to dynamically update the PDF path
# def create_pdf_tool(file_path):
#     if not os.path.exists(file_path):
#         raise FileNotFoundError(f"File {file_path} does not exist.")
#     return PDFSearchTool(pdf=file_path)





class FileProcessor:
    def __init__(self):
        self.supported_formats = {
            # Document formats
            '.pdf': self.extract_from_pdf,
            '.docx': self.extract_from_docx,
            '.txt': self.extract_from_txt,
            '.md': self.extract_from_markdown,
            
            # Spreadsheet formats
            '.xlsx': self.extract_from_excel,
            '.xls': self.extract_from_excel,
            '.csv': self.extract_from_csv,
            
            # Presentation formats
            '.pptx': self.extract_from_powerpoint,
            '.ppt': self.extract_from_powerpoint,
            
            # Audio formats
            '.mp3': self.extract_from_audio,
            '.wav': self.extract_from_audio,
            '.m4a': self.extract_from_audio,
            
            # Video formats
            '.mp4': self.extract_from_video,
            '.mov': self.extract_from_video,
            '.avi': self.extract_from_video,
            
            # Web formats
            '.json': self.extract_from_json,
            '.html': self.extract_from_html,
        }

    def extract_text_from_file(self, file_path: str) -> str:
        """Main function to extract text from various file formats"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_extension = Path(file_path).suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        try:
            return self.supported_formats[file_extension](file_path)
        except Exception as e:
            raise Exception(f"Error extracting text from {file_path}: {str(e)}")

    def extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            reader = PdfReader(file_path)
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"PDF extraction error: {str(e)}")

    def extract_from_excel(self, file_path: str) -> str:
        """Extract text from Excel file"""
        try:
            # Read all sheets
            excel_data = []
            df = pd.read_excel(file_path, sheet_name=None)
            
            for sheet_name, sheet_data in df.items():
                excel_data.append(f"\nSheet: {sheet_name}")
                
                # Convert headers and data to string
                headers = sheet_data.columns.tolist()
                excel_data.append(f"Headers: {', '.join(str(h) for h in headers)}")
                
                # Convert data to string format
                for _, row in sheet_data.iterrows():
                    row_data = [str(cell) for cell in row]
                    excel_data.append(" | ".join(row_data))
            
            return "\n".join(excel_data)
        except Exception as e:
            raise Exception(f"Excel extraction error: {str(e)}")

    def extract_from_csv(self, file_path: str) -> str:
        """Extract text from CSV file"""
        try:
            df = pd.read_csv(file_path)
            csv_data = []
            
            # Convert headers and data to string
            headers = df.columns.tolist()
            csv_data.append(f"Headers: {', '.join(str(h) for h in headers)}")
            
            # Convert data to string format
            for _, row in df.iterrows():
                row_data = [str(cell) for cell in row]
                csv_data.append(" | ".join(row_data))
            
            return "\n".join(csv_data)
        except Exception as e:
            raise Exception(f"CSV extraction error: {str(e)}")

    def extract_from_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text_runs = []
            
            for i, slide in enumerate(prs.slides):
                text_runs.append(f"\nSlide {i+1}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            
            return "\n".join(text_runs)
        except Exception as e:
            raise Exception(f"PowerPoint extraction error: {str(e)}")

    def extract_from_audio(self, file_path: str) -> str:
        """Extract text from audio file using speech recognition"""
        try:
            # Initialize recognizer
            recognizer = sr.Recognizer()
            
            # Convert audio to WAV if it's not already
            if not file_path.lower().endswith('.wav'):
                from pydub import AudioSegment
                audio = AudioSegment.from_file(file_path)
                wav_path = file_path + '.wav'
                audio.export(wav_path, format='wav')
                file_path = wav_path
            
            # Perform speech recognition
            with sr.AudioFile(file_path) as source:
                audio = recognizer.record(source)
                text = recognizer.recognize_google(audio)
                
            # Clean up temporary WAV file if created
            if file_path.endswith('.wav') and os.path.exists(file_path):
                os.remove(file_path)
                
            return text
        except Exception as e:
            raise Exception(f"Audio extraction error: {str(e)}")

    def extract_from_video(self, file_path: str) -> str:
        """Extract text from video file by converting audio and using speech recognition"""
        try:
            # Extract audio from video
            video = VideoFileClip(file_path)
            audio = video.audio
            
            # Save audio temporarily
            temp_audio = file_path + '.wav'
            audio.write_audiofile(temp_audio)
            
            # Extract text from audio
            text = self.extract_from_audio(temp_audio)
            
            # Cleanup
            video.close()
            audio.close()
            if os.path.exists(temp_audio):
                os.remove(temp_audio)
                
            return text
        except Exception as e:
            raise Exception(f"Video extraction error: {str(e)}")

    def extract_from_json(self, file_path: str) -> str:
        """Extract text from JSON file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                return json.dumps(data, indent=2)
        except Exception as e:
            raise Exception(f"JSON extraction error: {str(e)}")

    def extract_from_html(self, file_path: str) -> str:
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                return soup.get_text(separator='\n', strip=True)
        except Exception as e:
            raise Exception(f"HTML extraction error: {str(e)}")

    def extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            text = docx2txt.process(file_path)
            return text
        except Exception as e:
            raise Exception(f"DOCX extraction error: {str(e)}")

    def extract_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            raise Exception(f"TXT extraction error: {str(e)}")

    def extract_from_markdown(self, file_path: str) -> str:
        """Extract text from Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
                html_content = markdown.markdown(md_content)
                return re.sub('<[^<]+?>', '', html_content)
        except Exception as e:
            raise Exception(f"Markdown extraction error: {str(e)}")

def generate_unique_content(content: str, week: int, day: str, platform: str) -> str:
    """
    Generate unique content based on week and day context
    """
    # Split content into sentences
    sentences = content.split('. ')
    
    # Use week and day to select different portions of content
    start_idx = ((week - 1) * 5 + ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"].index(day)) % len(sentences)
    
    # Take a different subset of sentences for each day
    num_sentences = min(5, len(sentences))  # Take up to 5 sentences
    selected_sentences = sentences[start_idx:start_idx + num_sentences]
    
    # If we need more sentences, wrap around to the beginning
    if len(selected_sentences) < num_sentences:
        selected_sentences.extend(sentences[:num_sentences - len(selected_sentences)])
    
    # Join sentences back together
    unique_content = '. '.join(selected_sentences)
    
    # Add platform-specific formatting
    if platform == "twitter":
        unique_content = f"{unique_content}"
    elif platform == "instagram":
        unique_content = f"{unique_content}"
    elif platform == "linkedin":
        unique_content = f" {unique_content}"
    elif platform == "facebook":
        unique_content = f" {unique_content}"
    
    return unique_content


def generate_different_content(content, week, day, platform, post_number):
    """
    Generate unique content based on week, day context, and post number to ensure
    different content for multiple posts on the same day
    """
    # Split content into sentences
    sentences = content.split('. ')
    
    # Calculate different starting points for different posts on the same day
    base_idx = ((week - 1) * 5 + ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"].index(day)) % len(sentences)
    post_offset = (post_number - 1) * 3  # Offset by 3 sentences for each post
    start_idx = (base_idx + post_offset) % len(sentences)
    
    # Take a different subset of sentences for each post
    num_sentences = min(15, len(sentences))  # Take up to 5 sentences
    selected_sentences = sentences[start_idx:start_idx + num_sentences]
    
    # If we need more sentences, wrap around to the beginning
    if len(selected_sentences) < num_sentences:
        remaining = num_sentences - len(selected_sentences)
        selected_sentences.extend(sentences[:remaining])
    
    # Join sentences back together
    unique_content = '. '.join(selected_sentences)
    
    # Add platform-specific formatting with post number
    if platform == "twitter":
        unique_content = f" {unique_content}"
    elif platform == "instagram":
        unique_content = f"{unique_content}"
    elif platform == "linkedin":
        unique_content = f"{unique_content}"
    elif platform == "facebook":
        unique_content = f"{unique_content}"
    elif platform == "youtube":
        unique_content == f" {unique_content}"
    elif platform == "tiktok":
        unique_content == f" {unique_content}"
    
    return unique_content





def extract_title_from_content(content: str) -> str:
    """
    Extract or generate a title from the content
    """
    # Try to get the first sentence
    sentences = content.split('. ')
    if sentences:
        # Use the first sentence if it's not too long
        first_sentence = sentences[0].strip()
        if len(first_sentence) <= 50:
            return first_sentence
        
        # If first sentence is too long, take first few words
        words = first_sentence.split()
        title = ' '.join(words[:7])
        return f"{title}..."
    
    return "Untitled Post"

def process_content_for_platform(content: str, platform: str, limits: Dict) -> str:
    """
    Process and format content according to platform-specific limits
    """
    # Clean content
    content = content.strip()
    
    # Apply character limit if specified
    if limits.get("chars"):
        content = content[:limits["chars"]]
    
    # Apply word limit if specified
    if limits.get("words"):
        words = content.split()
        content = ' '.join(words[:limits["words"]])
    
    # Add platform-specific formatting
    if platform == "twitter":
        # Leave room for hashtags
        
        content += ""
    
    elif platform == "instagram":
        content += ""
    
    elif platform == "linkedin":
        content += ""
    
    elif platform == "facebook":
        content += ""

    elif platform == "youtube":
        content += ""

    elif platform == "tiktok":
        content == ""
    
    return content.strip()


# Instantiate FileProcessor
file_processor = FileProcessor()

# Define a tool wrapper for file processing
# file_extraction_tool = Tool(
#     name="FileExtractionTool",
#     description="A universal tool for extracting text from various file formats.",
#     func=lambda file_path: file_processor.extract_text_from_file(file_path),  # Use 'func' as the field name
# )



def trim_content(content: str, platform: str, limits: dict) -> str:
    """Trim content based on platform-specific character or word limits."""
    char_limit = limits.get(platform, {}).get("chars")
    word_limit = limits.get(platform, {}).get("words")
    
    # Apply character limit
    if char_limit:
        content = content[:char_limit].strip()
    
    # Apply word limit
    if word_limit:
        words = content.split()
        content = " ".join(words[:word_limit]).strip()
    
    return content



def generate_twitter_post(content: str, limits: dict) -> str:
    """Generate a concise Twitter post with dynamic hashtags."""
    content = trim_content(content, "twitter", limits)
    hashtags = ""
    return f"{content}\n\n{hashtags}"

# twitter_post_tool = Tool(
#     name="TwitterPostGeneratorTool",
#     description="Generates concise and impactful tweets with dynamic hashtags, respecting platform limits.",
#     func=lambda content: generate_twitter_post(content, PLATFORM_LIMITS),
# )



def generate_instagram_post(content: str, limits: dict) -> str:
    """Generate an Instagram post with emojis and hashtags."""
    content = trim_content(content, "instagram", limits)
    hashtags = "#InstaLife #Photography #Explore"
    emoji = ""
    return f"{emoji} {content.strip()} {hashtags}"

# instagram_post_tool = Tool(
#     name="InstagramPostGeneratorTool",
#     description="Generates Instagram captions with emojis and hashtags, respecting platform limits.",
#     func=lambda content: generate_instagram_post(content, PLATFORM_LIMITS),
# )


def generate_linkedin_post(content: str, limits: dict) -> str:
    """Generate a professional LinkedIn post with hashtags."""
    content = trim_content(content, "linkedin", limits)
    hashtags = ""
    return f"{content.strip()} {hashtags} Engage with this post by sharing your thoughts!"

# linkedin_post_tool = Tool(
#     name="LinkedInPostGeneratorTool",
#     description="Generates professional LinkedIn posts with dynamic formatting and hashtags.",
#     func=lambda content: generate_linkedin_post(content, PLATFORM_LIMITS),
# )


def generate_facebook_post(content: str, limits: dict) -> str:
    """Generate a friendly Facebook post."""
    content = trim_content(content, "facebook", limits)
    return f"{content.strip()} "

# facebook_post_tool = Tool(
#     name="FacebookPostGeneratorTool",
#     description="Generates friendly and engaging Facebook posts, respecting platform limits.",
#     func=lambda content: generate_facebook_post(content, PLATFORM_LIMITS),
# )


def generate_wordpress_post(content: str, limits: dict) -> str:
    """Generate a WordPress blog post with SEO-friendly structure."""
    content = trim_content(content, "wordpress", limits)
    # title = "Blog Title: " + (content[:50] + "..." if len(content) > 50 else content)
    return f" {content.strip()} #SEO #WordPress"

# wordpress_post_tool = Tool(
#     name="WordPressPostGeneratorTool",
#     description="Generates SEO-optimized WordPress blog posts, respecting platform limits.",
#     func=lambda content: generate_wordpress_post(content, PLATFORM_LIMITS),
# )


def generate_youtube_post(content: str, limits: dict) -> str:
    """Generate a youtube post which is further used for image or video genration"""
    content = trim_content(content, "youtube", limits)
    # title = "Blog Title: " + (content[:50] + "..." if len(content) > 50 else content)
    return f" {content.strip()} "

# youtube_post_tool = Tool(
#     name="youtubePostGeneratorTool",
#     description="Generate a youtube post which is further used for image or video genration.",
#     func=lambda content: generate_youtube_post(content, PLATFORM_LIMITS),
# )


def generate_tiktok_post(content: str, limits: dict) -> str:
    """Generate a tiktok post which is further used for image or video genration"""
    content = trim_content(content, "tiktok", limits)
    # title = "Blog Title: " + (content[:50] + "..." if len(content) > 50 else content)
    return f" {content.strip()} "

# tiktok_post_tool = Tool(
#     name="tiktokPostGeneratorTool",
#     description="Generate a youtube post which is further used for image or video genration.",
#     func=lambda content: generate_tiktok_post(content, PLATFORM_LIMITS),
# )



# Headers for Airtable API
HEADERS = {
    "Authorization": f"Bearer {AIRTABLE_API_TOKEN}"
}

def fetch_airtable_records():
    """Fetch records from Airtable."""
    response = requests.get(AIRTABLE_URL, headers=HEADERS)
    if response.status_code == 200:
        return response.json().get("records", [])
    return []

def find_matching_prompt(user_prompt, records):
    """Find the best matching DALL·E prompt from Airtable based on the query."""
    for record in records:
        fields = record.get("fields", {})
        keywords = fields.get("dalle keyword", "").split(", ")  # Comma-separated keywords
        
        # Check if any keyword is in the user prompt
        if any(keyword.lower() in user_prompt.lower() for keyword in keywords):
            return fields.get("dalle prompt")
    return None  # No match found

def upload_image_to_sftp(image_data: bytes, filename: str) -> str:
    """Upload image data to SFTP server and return the permanent URL."""
    try:
        if not all([SFTP_HOST, SFTP_USER, SFTP_PASS]):
            logger.warning("SFTP credentials not configured. Returning temporary DALL-E URL.")
            return None
        
        remote_path = f"/home/{SFTP_USER}/public_html/nishant/{filename}"
        base_url = f"https://{SFTP_HOST.replace('sftp.', '').replace('ftp.', '')}/nishant/{filename}"
        
        # Connect to SFTP and upload the file
        ssh_client = paramiko.SSHClient()
        ssh_client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        ssh_client.connect(hostname=SFTP_HOST, port=SFTP_PORT, username=SFTP_USER, password=SFTP_PASS)
        sftp = ssh_client.open_sftp()
        
        # Upload the file
        with sftp.file(remote_path, 'wb') as remote_file:
            remote_file.write(image_data)
        
        sftp.close()
        ssh_client.close()
        
        logger.info(f"✅ Uploaded image to SFTP: {base_url}")
        return base_url
        
    except paramiko.AuthenticationException:
        logger.error("SFTP authentication failed. Check credentials.")
        return None
    except paramiko.SSHException as ssh_e:
        logger.error(f"SFTP error: {str(ssh_e)}")
        return None
    except Exception as e:
        logger.error(f"Failed to upload image to SFTP: {str(e)}")
        return None

def generate_image(query, content):
    """Generate an image using OpenAI's DALL·E, download it, and upload to permanent storage."""
    records = fetch_airtable_records()
    style_prompt = find_matching_prompt(query, records)
    
    # Combine content with style prompt if found, otherwise use content alone
    final_prompt = f"{content} {style_prompt}" if style_prompt else content

    try:
        # Generate image with DALL·E
        response = client.images.generate(
            model="dall-e-3",
            prompt=final_prompt,
            size="1024x1024",
            n=1
        )
        
        dall_e_url = response.data[0].url
        logger.info(f"🖼️ Generated DALL·E image: {dall_e_url[:100]}...")
        
        # CRITICAL: Download the image immediately (whether it's DALL-E URL or Azure blob URL)
        # Azure blob URLs expire, so we MUST download and store permanently
        image_data = None
        try:
            img_response = requests.get(dall_e_url, timeout=60)  # Increased timeout for large images
            img_response.raise_for_status()
            image_data = img_response.content
            logger.info(f"📥 Downloaded image from source ({len(image_data)} bytes)")
        except Exception as download_error:
            logger.error(f"❌ CRITICAL: Failed to download image from source: {download_error}")
            logger.error(f"❌ Image URL was: {dall_e_url[:200]}...")
            # DO NOT return temporary URL - raise error instead
            raise Exception(f"Failed to download image: {str(download_error)}. Image must be downloaded immediately to prevent expiration.")
        
        if not image_data:
            raise Exception("Downloaded image data is empty. Cannot proceed with permanent storage.")
        
        # Generate a unique filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        import hashlib
        content_hash = hashlib.md5(image_data).hexdigest()[:8]
        filename = f"dalle_{timestamp}_{content_hash}.png"
        
        # Upload to permanent storage (SFTP) - REQUIRED, no fallback
        permanent_url = upload_image_to_sftp(image_data, filename)
        
        if permanent_url:
            logger.info(f"✅ Image stored permanently: {permanent_url}")
            return permanent_url
        else:
            # CRITICAL: SFTP upload failed - log error but don't return temporary URL
            logger.error(f"❌ CRITICAL: SFTP upload failed. Image downloaded but not stored permanently.")
            logger.error(f"❌ SFTP credentials check: HOST={bool(SFTP_HOST)}, USER={bool(SFTP_USER)}, PASS={bool(SFTP_PASS)}")
            # Raise error instead of returning temporary URL
            raise Exception("SFTP upload failed. Image was downloaded but could not be stored permanently. Please check SFTP configuration.")
            
    except Exception as e:
        raise Exception(f"Error generating image: {str(e)}")
    


# # FTP credentials (load from environment variables)
# FTP_HOST = os.getenv("ftp_host")
# FTP_USER = os.getenv("ftp_user")
# FTP_PASS = os.getenv("ftp_pass")

# def delete_image_from_ftp(image_url: str):
#     """Delete an image from the FTP server given its URL."""
#     try:
#         # Parse the URL to extract the filename
#         parsed_url = urlparse(image_url)
#         filename = os.path.basename(parsed_url.path)
        
#         # Connect to FTP and delete the file
#         with FTP(FTP_HOST) as ftp:
#             ftp.login(user=FTP_USER, passwd=FTP_PASS)
#             ftp.delete(filename)
#             logger.info(f"Deleted image {filename} from FTP")
#     except Exception as e:
#         logger.error(f"Failed to delete image from FTP: {str(e)}")


# SFTP credentials (load from environment variables)
SFTP_HOST = os.getenv("SFTP_HOST")
SFTP_USER = os.getenv("SFTP_USER")
SFTP_PASS = os.getenv("SFTP_PASS")
SFTP_PORT = int(os.getenv("SFTP_PORT", "22"))  # Default to 22 if not specified

def delete_image_from_ftp(image_url: str):
    """Delete an image from the SFTP server given its URL."""
    try:
        # Parse the URL to extract the filename
        parsed_url = urlparse(image_url)
        filename = os.path.basename(parsed_url.path)
        remote_path = f"/home/{SFTP_USER}/public_html/nishant/{filename}"

        # Connect to SFTP and delete the file
        ssh_client = paramiko.SSHClient()
        ssh_client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        ssh_client.connect(hostname=SFTP_HOST, port=SFTP_PORT, username=SFTP_USER, password=SFTP_PASS)
        sftp = ssh_client.open_sftp()

        # Delete the file
        sftp.remove(remote_path)
        sftp.close()
        ssh_client.close()
        logger.info(f"Deleted image {filename} from SFTP at {remote_path}")

    except paramiko.AuthenticationException:
        logger.error("SFTP authentication failed. Check credentials.")
        raise Exception("SFTP authentication failed. Check credentials.")
    except paramiko.SSHException as ssh_e:
        logger.error(f"SFTP error: {str(ssh_e)}")
        raise Exception(f"SFTP error: {str(ssh_e)}")
    except FileNotFoundError:
        logger.error(f"Image not found on SFTP server: {remote_path}")
        raise Exception(f"Image not found on SFTP server: {filename}")
    except Exception as e:
        logger.error(f"Failed to delete image from SFTP: {str(e)}")
        raise Exception(f"Failed to delete image from SFTP: {str(e)}")